import re
import json
import argparse
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from io import BytesIO
from datetime import datetime

import pandas as pd
from pptx import Presentation

from PIL import Image
import imagehash


# -----------------------------
# Config
# -----------------------------
REF_SECTION_HINTS = [
    r"레퍼런스",
    r"reference",
    r"사례",
    r"운영\s*사례",
    r"주요\s*진행\s*사례",
    r"유사\s*교육",
    r"교육\s*레퍼런스",
]

FIELD_HINTS = [
    r"교육\s*개요",
    r"수강\s*대상",
    r"교육\s*대상",
    r"교육\s*형태",
    r"교육\s*방식",
    r"교육\s*시수",
    r"과정\s*구성",
    r"교육\s*주제",
    r"교육\s*목적",
    r"기업\s*니즈",
]

EXCLUDE_TITLE_HINTS = [
    r"레퍼런스",
    r"reference",
    r"사례",
    r"유사\s*교육",
    r"주요\s*진행\s*사례",
    r"교육\s*레퍼런스",
    r"교육\s*개요",
    r"과정\s*구성",
]

HEADER_TO_KEY = {
    "기업의 니즈": "기업 니즈",
    "기업 니즈": "기업 니즈",
    "교육 목적": "교육 목적",
    "교육 주제": "교육 주제",
    "교육 대상": "교육 대상",
    "수강 대상": "수강 대상",
    "교육 형태": "교육 형태",
    "교육 방식": "교육 방식",
    "교육 시수": "교육 시수",
}


# -----------------------------
# Utils: logo hash map
# -----------------------------
def load_logo_map(path: Path) -> Dict[str, str]:
    if path.exists():
        try:
            return json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}

def save_logo_map(path: Path, logo_map: Dict[str, str]) -> None:
    path.write_text(json.dumps(logo_map, ensure_ascii=False, indent=2), encoding="utf-8")

def compute_phash_from_blob(blob: bytes) -> Optional[str]:
    try:
        img = Image.open(BytesIO(blob))
        # 팔레트/투명도 경고 방지
        if img.mode in ("P", "LA"):
            img = img.convert("RGBA")
        else:
            img = img.convert("RGB")
        return str(imagehash.phash(img))
    except Exception:
        return None

def extract_slide_logo_hashes(slide) -> List[str]:
    hashes: List[str] = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                blob = shape.image.blob
                h = compute_phash_from_blob(blob)
                if h:
                    hashes.append(h)
            except Exception:
                pass
    return hashes


# -----------------------------
# PPTX -> slide blocks/lines + logo hashes
# -----------------------------
def extract_slides(pptx_path: Path) -> List[Dict[str, Any]]:
    prs = Presentation(str(pptx_path))
    out: List[Dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        blocks: List[Dict[str, Any]] = []

        for shape in slide.shapes:
            top = int(getattr(shape, "top", 0))

            # 1) table
            if getattr(shape, "has_table", False):
                tbl = shape.table
                for r in range(len(tbl.rows)):
                    for c in range(len(tbl.columns)):
                        cell = tbl.cell(r, c)
                        txt = (cell.text or "").strip()
                        if txt:
                            for t in re.split(r"[\r\n]+", txt):
                                t = t.strip()
                                if t:
                                    blocks.append({
                                        "top": top,
                                        "text": t,
                                        "level": 0,
                                        "kind": "table",
                                    })
                continue

            # 2) text frame
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                tf = shape.text_frame
                for p in tf.paragraphs:
                    txt = (p.text or "").strip()
                    if not txt:
                        continue
                    lvl = int(getattr(p, "level", 0))
                    for t in re.split(r"[\r\n]+", txt):
                        t = t.strip()
                        if t:
                            blocks.append({
                                "top": top,
                                "text": t,
                                "level": lvl,
                                "kind": "text",
                            })

        blocks.sort(key=lambda b: (b["top"], b["level"], b["text"]))
        lines = [b["text"] for b in blocks]
        logo_hashes = extract_slide_logo_hashes(slide)

        out.append({
            "file": pptx_path.name,
            "slide_index": idx,
            "blocks": blocks,
            "lines": lines,
            "full_text": "\n".join(lines),
            "logo_hashes": logo_hashes,
        })

    return out


# -----------------------------
# Candidate detection
# -----------------------------
def score_reference_candidate(slide: Dict[str, Any]) -> Tuple[bool, int, List[str]]:
    text = slide["full_text"].lower()
    hits: List[str] = []
    score = 0

    for pat in REF_SECTION_HINTS:
        if re.search(pat, text, flags=re.IGNORECASE):
            score += 2
            hits.append(pat)

    for pat in FIELD_HINTS:
        if re.search(pat, text, flags=re.IGNORECASE):
            score += 1
            hits.append(pat)

    return (score >= 3), score, sorted(set(hits))


# -----------------------------
# Parsing helpers
# -----------------------------
def parse_kv_line(line: str) -> Optional[Tuple[str, str]]:
    """
    라벨형 추출:
    - "교육 주제 | xxx"
    - "교육 대상: xxx"
    - "교육 시수 - 총 7시간"
    """
    s = line.strip()

    keys = r"(교육\s*주제|교육\s*대상|수강\s*대상|교육\s*형태|교육\s*방식|교육\s*시수|교육\s*목적|기업\s*니즈|기업의\s*니즈)"
    m = re.match(rf"^{keys}\s*\|\s*(.+)$", s)
    if m:
        return m.group(1).strip(), m.group(2).strip()

    m = re.match(rf"^{keys}\s*[:：]\s*(.+)$", s)
    if m:
        return m.group(1).strip(), m.group(2).strip()

    m = re.match(rf"^{keys}\s*[-–]\s*(.+)$", s)
    if m:
        return m.group(1).strip(), m.group(2).strip()

    return None


def normalize_duration_hours(raw: str) -> Optional[float]:
    if not raw:
        return None
    s = raw.lower().replace(" ", "")

    m = re.search(r"(\d+(?:\.\d+)?)h?x(\d+(?:\.\d+)?)", s)
    if m:
        try:
            return float(m.group(1)) * float(m.group(2))
        except Exception:
            pass

    m = re.search(r"(\d+(?:\.\d+)?)(?:시간|h)", s)
    if m:
        try:
            return float(m.group(1))
        except Exception:
            return None

    return None


def guess_title(lines: List[str]) -> str:
    for ln in lines[:8]:
        if any(re.search(p, ln, flags=re.IGNORECASE) for p in EXCLUDE_TITLE_HINTS):
            continue
        if any(re.search(p, ln, flags=re.IGNORECASE) for p in FIELD_HINTS):
            continue
        if len(ln.strip()) >= 6:
            return ln.strip()
    return ""


def enrich_kv_with_header_nextline(lines: List[str], kv: Dict[str, str]) -> None:
    """
    예) "기업의 니즈" 다음 줄이 실제 내용인 경우를 보강.
    이미 kv에 값이 있으면 덮어쓰지 않음.
    """
    for i, ln in enumerate(lines[:-1]):
        s = ln.strip()
        nxt = lines[i + 1].strip()

        if s in HEADER_TO_KEY and parse_kv_line(s) is None:
            key = HEADER_TO_KEY[s]

            # 다음 줄이 또다른 헤더/라벨이면 skip
            if nxt in HEADER_TO_KEY:
                continue
            if parse_kv_line(nxt) is not None:
                continue
            if len(nxt) >= 3:
                kv.setdefault(key, nxt)


def build_details(slide: Dict[str, Any]) -> List[str]:
    blocks = slide.get("blocks", [])
    details: List[str] = []

    def is_header(text: str) -> bool:
        # 과정 구성 / [과정 구성] / ＜과정 구성＞ 등 허용
        return bool(re.fullmatch(r"\s*[\[\(<{＜]?\s*과정\s*구성\s*[\]\)>}＞]?\s*", text))

    def looks_like_section_break(text: str) -> bool:
        # "교안/예시" 같은 확실한 종료만
        return bool(re.search(r"(교안\s*자료|부록|참고|예시|실습\s*예시|<.*예시>|＜.*예시＞)", text))

    start_idx = None
    for i, b in enumerate(blocks):
        if is_header(b["text"]):
            start_idx = i
            break
    if start_idx is None:
        return []

    # 헤더 바로 위에 트랙명/과정명 같은 짧은 라인이 있는 경우 1개 보강
    for j in range(start_idx - 1, max(-1, start_idx - 5), -1):
        tt = blocks[j]["text"].strip()
        if not tt:
            continue
        if parse_kv_line(tt) is not None:
            continue
        if re.search(r"(레퍼런스|사례|교육\s*레퍼런스)", tt, re.IGNORECASE):
            continue
        if len(tt) <= 40:
            details.append(tt)
            break

    # 헤더 아래 수집
    for b in blocks[start_idx + 1:]:
        t = b["text"].strip()
        if not t:
            continue

        if looks_like_section_break(t):
            break

        # KV 라인은 details에 넣지 않음(하지만 종료도 아님)
        if parse_kv_line(t) is not None:
            continue

        # bullet
        m = re.match(r"^[-•·]\s*(.+)$", t)
        if m:
            item = m.group(1).strip()
            if item:
                details.append(item)
            continue

        # level-based bullets
        if int(b.get("level", 0)) > 0:
            details.append(t)
            continue

        # table cells
        if b.get("kind") == "table":
            details.append(t)
            continue

        # 일반 라인도 구성으로 쓰일 수 있음
        if len(t) >= 2:
            details.append(t)

    # 후처리: 너무 명백한 메타 라인 제거
    cleaned = []
    for d in details:
        if re.search(r"(교육\s*대상|교육\s*시수|교육\s*주제|교육\s*형태|교육\s*방식)", d):
            continue
        cleaned.append(d)

    # 중복 제거(순서 유지)
    seen = set()
    uniq = []
    for d in cleaned:
        if d not in seen:
            seen.add(d)
            uniq.append(d)
    return uniq


def parse_ref_case(slide: Dict[str, Any], logo_map: Dict[str, str]) -> Tuple[Optional[Dict[str, Any]], List[Dict[str, Any]]]:
    """
    return: (case or None, issues[])
    """
    issues: List[Dict[str, Any]] = []
    lines = slide["lines"]
    kv: Dict[str, str] = {}

    # 1) KV 라인 파싱
    for ln in lines:
        kv_pair = parse_kv_line(ln)
        if kv_pair:
            k, v = kv_pair
            # normalize key: "기업의 니즈" -> "기업 니즈"
            k_norm = HEADER_TO_KEY.get(k.strip(), k.strip())
            kv[k_norm] = v.strip()

    # 2) 헤더-다음줄 보강
    enrich_kv_with_header_nextline(lines, kv)

    # 3) 최소 기준 체크
    core_keys = ["교육 주제", "교육 대상", "수강 대상", "교육 형태", "교육 방식", "교육 시수", "기업 니즈", "교육 목적"]
    found = sum(1 for k in core_keys if k in kv and str(kv.get(k, "")).strip())
    if found < 2:
        return None, issues

    topic = kv.get("교육 주제", "") or ""
    target = kv.get("수강 대상", "") or kv.get("교육 대상", "") or ""
    fmt = kv.get("교육 형태", "") or kv.get("교육 방식", "") or ""
    duration_raw = kv.get("교육 시수", "") or ""
    duration = normalize_duration_hours(duration_raw)

    title = guess_title(lines)
    details = build_details(slide)

    # client from logo map
    client = ""
    for h in slide.get("logo_hashes", []):
        mapped = (logo_map.get(h) or "").strip()
        if mapped:
            client = mapped
            break
    if not client and slide.get("logo_hashes"):
        client = f"UNKNOWN_LOGO:{slide['logo_hashes'][0]}"

    # needs: 기업 니즈 우선, 없으면 교육 목적
    needs = (kv.get("기업 니즈", "") or kv.get("교육 목적", "") or "").strip()

    # issues 기록
    preview = " / ".join([x.strip() for x in lines[:10] if x.strip()])[:400]
    if not needs:
        issues.append({
            "issue_type": "missing_needs",
            "file": slide["file"],
            "slide_index": slide["slide_index"],
            "message": "needs(기업의 니즈/교육 목적)를 추출하지 못함",
            "preview": preview,
        })
    if not details:
        issues.append({
            "issue_type": "missing_details",
            "file": slide["file"],
            "slide_index": slide["slide_index"],
            "message": "details(과정 구성)를 추출하지 못함(헤더/표/텍스트 형태 확인 필요)",
            "preview": preview,
        })
    if client.startswith("UNKNOWN_LOGO:") or not client:
        issues.append({
            "issue_type": "unknown_client_logo",
            "file": slide["file"],
            "slide_index": slide["slide_index"],
            "message": "client가 로고 이미지 기반 UNKNOWN 상태 (logo_hash_map.json 매핑 필요)",
            "preview": preview,
        })

    # confidence
    score = 0.0
    if title: score += 0.2
    if topic: score += 0.2
    if target: score += 0.2
    if fmt: score += 0.2
    if duration is not None: score += 0.2
    confidence = round(score, 2)

    case = {
        "client": client,
        "needs": needs,
        "title": title,
        "target": target,
        "format": fmt,
        "topic": topic,
        "duration": duration,
        "details": details,
        "source_file": slide["file"],
        "source_slide_index": slide["slide_index"],
        "confidence_score": confidence,
    }

    return case, issues


# -----------------------------
# Master append helpers
# -----------------------------
def load_existing_keys(master_csv: Path) -> set:
    """
    중복 방지: (source_file, source_slide_index) 키로 이미 누적된 항목은 skip
    """
    if not master_csv.exists():
        return set()
    try:
        df = pd.read_csv(master_csv, dtype={"source_slide_index": int})
        keys = set(zip(df["source_file"].astype(str), df["source_slide_index"].astype(int)))
        return keys
    except Exception:
        return set()

def append_csv(master_csv: Path, rows: List[Dict[str, Any]], columns: List[str]) -> int:
    if not rows:
        return 0
    df = pd.DataFrame(rows)
    # 컬럼 정렬/누락 보정
    for c in columns:
        if c not in df.columns:
            df[c] = ""
    df = df[columns]
    header = not master_csv.exists()
    df.to_csv(master_csv, mode="a", index=False, header=header, encoding="utf-8-sig")
    return len(df)


# -----------------------------
# Runner
# -----------------------------
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", required=True, help="pptx file or folder containing pptx files")
    ap.add_argument("--out", default="out", help="output folder (master files live here)")
    ap.add_argument("--append", action="store_true", help="append to references_master.csv / issues_master.csv (dedupe by file+slide)")
    args = ap.parse_args()

    in_path = Path(args.input)
    out_dir = Path(args.out)
    out_dir.mkdir(parents=True, exist_ok=True)

    # input list
    if in_path.is_dir():
        pptx_files = sorted(in_path.glob("*.pptx"))
    else:
        pptx_files = [in_path]

    # master paths
    references_master = out_dir / "references_master.csv"
    issues_master = out_dir / "issues_master.csv"
    logo_map_path = out_dir / "logo_hash_map.json"

    # dedupe keys when append
    existing_keys = load_existing_keys(references_master) if args.append else set()

    logo_map = load_logo_map(logo_map_path)

    total_slides = 0
    candidate_count = 0
    parsed_count = 0
    skipped_dup = 0
    new_issues = 0

    run_id = datetime.now().strftime("%Y%m%d_%H%M%S")

    new_refs: List[Dict[str, Any]] = []
    new_issues_rows: List[Dict[str, Any]] = []

    for f in pptx_files:
        slides = extract_slides(f)
        total_slides += len(slides)

        for s in slides:
            ok, _, _ = score_reference_candidate(s)
            if not ok:
                continue
            candidate_count += 1

            key = (s["file"], int(s["slide_index"]))
            if args.append and key in existing_keys:
                skipped_dup += 1
                continue

            case, issues = parse_ref_case(s, logo_map)
            if case:
                parsed_count += 1
                case["run_id"] = run_id
                # details는 CSV에 저장할 때 합쳐서 넣기 좋게 문자열 컬럼도 함께
                case["details_joined"] = " | ".join(case["details"]) if isinstance(case.get("details"), list) else str(case.get("details") or "")
                new_refs.append(case)

                for it in issues:
                    it["run_id"] = run_id
                    new_issues_rows.append(it)

    # write/append masters
    ref_cols = [
        "run_id",
        "client", "needs", "title", "target", "format", "topic", "duration",
        "details_joined",
        "source_file", "source_slide_index",
        "confidence_score",
    ]
    issue_cols = ["run_id", "issue_type", "file", "slide_index", "message", "preview"]

    written_refs = append_csv(references_master, new_refs, ref_cols) if (args.append or not references_master.exists()) else 0
    written_issues = append_csv(issues_master, new_issues_rows, issue_cols) if (args.append or not issues_master.exists()) else 0

    save_logo_map(logo_map_path, logo_map)

    # summary print only
    print("[DONE] Master outputs:")
    print(f"- {references_master}")
    print(f"- {issues_master}")
    print(f"- {logo_map_path}")
    print("[STATS]")
    print(f"- input_files: {len(pptx_files)}")
    print(f"- total_slides: {total_slides}")
    print(f"- candidates: {candidate_count}")
    print(f"- parsed_refs: {parsed_count}")
    if args.append:
        print(f"- skipped_duplicates: {skipped_dup}")
    print(f"- written_refs: {written_refs}")
    print(f"- written_issues: {written_issues}")
    print(f"- run_id: {run_id}")


if __name__ == "__main__":
    main()
