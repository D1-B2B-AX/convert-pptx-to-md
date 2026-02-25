import os
import re
import json 
import pandas as pd
from pptx import Presentation
from openai import OpenAI
from dotenv import load_dotenv
from utils.pptx_parser import normalize, is_slide_hidden, get_visual_title, extract_text_from_slide

# =========================================================
# [설정] API 키 및 경로
# =========================================================
load_dotenv()
client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))

SOURCE_DIR = './input'
OUTPUT_DIR = './output/curriculum/'

# =========================================================
# [키워드 설정]
# =========================================================
EXCLUDE_KEYWORDS = [
    "유사", "사례", "실적", "reference", "case", "history", "result",
    "강사프로필", "수행실적", "제안사", "회사소개"
]

OVERVIEW_KEYWORDS = [
    "과정 소개", "과정소개", "과정 개요", "과정개요", 
    "교육 소개", "교육소개", "교육 개요", "교육개요", 
    "개요", "소개", "overview", "summary", "요약", "제안 배경", "기획 의도",
    "목표", "대상" 
]

CURRICULUM_KEYWORDS = [
    "커리큘럼", "세부과정", "교육과정", "교육내용", "모듈구성", 
    "상세과정", "프로그램", "module", "schedule", "curriculum",
    "모듈", "구성", "일정", "방법", "contents", "agenda", "syllabus",
    "1일차", "2일차", "1h", "2h", "time" 
]

# =========================================================
# [기능 1] PPTX 파싱 헬퍼 함수들
# =========================================================
def normalize(text):
    return re.sub(r'\s+', '', str(text).lower())

def get_visual_title(slide):
    if slide.shapes.title and slide.shapes.title.text.strip():
        return slide.shapes.title.text.strip()
    
    candidates = []
    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        # 상단에 위치한 텍스트를 제목 후보로 간주
        if shape.top < 2000000: 
            candidates.append((shape.top, shape.left, shape.text.strip()))
    
    if candidates:
        candidates.sort(key=lambda x: (x[0], x[1])) 
        return candidates[0][2]
    return ""

def check_table_headers(slide):
    for shape in slide.shapes:
        if shape.has_table:
            header_text = ""
            try:
                for cell in shape.table.rows[0].cells:
                    header_text += cell.text + " "
            except:
                continue
            norm_header = normalize(header_text)
            for key in CURRICULUM_KEYWORDS:
                if normalize(key) in norm_header:
                    return True
    return False

def classify_slide_advanced(slide):
    title = get_visual_title(slide)
    norm_title = normalize(title)
    
    for key in EXCLUDE_KEYWORDS:
        if normalize(key) in norm_title: return "EXCLUDE"
    for key in CURRICULUM_KEYWORDS:
        if normalize(key) in norm_title: return "CURRICULUM"
    for key in OVERVIEW_KEYWORDS:
        if normalize(key) in norm_title: return "OVERVIEW"
    if check_table_headers(slide):
        return "CURRICULUM"
    return "OTHER"

def extract_text_from_slide(slide):
    lines = []
    visual_title = get_visual_title(slide)
    if visual_title:
        lines.append(f"### {visual_title}")
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            if shape.text.strip() == visual_title:
                continue
            lines.append(shape.text.strip())
        
        if shape.has_table:
            for row in shape.table.rows:
                # 표 내용 한 줄로 합치기 (Markdown 변환 시 LLM이 처리하도록 원본 유지)
                row_cells = [c.text.replace('\n', ' ').strip() for c in row.cells if c.text.strip()]
                if row_cells:
                    lines.append(f"| {' | '.join(row_cells)} |")
    return "\n".join(lines)

# =========================================================
# [기능 2] LLM을 이용한 Markdown 변환
# =========================================================
def generate_rag_markdown(filename, course_idx, overview_text, curriculum_text):
    # 내용이 너무 적으면 스킵
    if len(curriculum_text) < 50: 
        return None

    # 요청하신 Metadata Block이 포함된 프롬프트
    prompt = f"""
    당신은 'B2B 교육 커리큘럼 정리 전문가'입니다.
    아래 제공된 Raw Text를 분석하여, RAG 검색에 최적화된 **Clean Markdown** 포맷으로 변환하십시오.

    [Input Source]
    - File: {filename}
    - Context: {overview_text[:3000]}
    - Content: {curriculum_text[:15000]}

    [Output Format Rules - Strict Markdown]
    1. **Metadata Block**: 문서 최상단에 아래 양식을 반드시 포함할 것.
       > **File**: {filename}
    
    2. **Section Structuring**:
       - 과정명/주제는 `# (H1)` 태그 사용
       - '교육 개요', '학습 목표' 등 대분류는 `## (H2)` 태그 사용
       - 세부 모듈/시간표는 `### (H3)` 태그 사용
    
    3. **Curriculum Table**:
       - 커리큘럼 상세 내용은 반드시 Markdown Table 혹은 계층형 List(`-`)로 정리할 것.
       - 시간(Time), 모듈명(Module), 세부내용(Detail)이 명확히 구분되어야 함.

    4. **Filtering**:
       - '강사 약력', '회사 홍보', '레퍼런스' 등 커리큘럼과 무관한 내용은 과감히 삭제할 것.
       - 정보가 없으면 없는 대로 놔둘 것 (지어내지 말 것).
       
    5. **No Chit-chat**: 서론/본론 없이 오직 Markdown 내용만 출력할 것. 만약 유효한 커리큘럼 정보가 없다면 오직 `NO_DATA`라고만 출력.
    """

    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            temperature=0
        )
        result = response.choices[0].message.content.strip()
        
        if "NO_DATA" in result: return None
        if len(result) < 50: return None
        
        return result

    except Exception as e:
        print(f"  ❌ LLM Error: {e}")
        return None

# =========================================================
# [기능 3] 메인 파이프라인 실행
# =========================================================
def process_rag_data_final():
    if not os.path.exists(SOURCE_DIR):
        print(f"❌ 원본 폴더를 찾을 수 없습니다: {SOURCE_DIR}")
        return

    # 폴더가 없으면 자동으로 생성 (output/curriculum)
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    files = [f for f in os.listdir(SOURCE_DIR) if f.endswith('.pptx')]
    print(f"🚀 총 {len(files)}개의 제안서 -> [커리큘럼] Markdown 변환 시작...\n")

    for file in files:
        file_path = os.path.join(SOURCE_DIR, file)
        print(f"📄 분석 중: {file}")

        try:
            prs = Presentation(file_path)
            courses = [] 
            current_course = {'overview': [], 'curriculum': []}

            # ---------------------------------------------------------
            # [핵심] PPTX 파싱 모듈을 활용한 깔끔한 반복문
            # ---------------------------------------------------------
            for i, slide in enumerate(prs.slides):
                
                # 1. 숨기기 처리된 슬라이드 통과
                if is_slide_hidden(slide):
                    continue

                # 2. 커리큘럼 추출용 슬라이드 분류
                slide_type = classify_slide_advanced(slide)
                if slide_type == "EXCLUDE": 
                    continue

                # 3. 텍스트 추출
                text = extract_text_from_slide(slide)

                # 4. 개요/커리큘럼 묶기 로직
                if slide_type == "OVERVIEW":
                    if current_course['curriculum']: 
                        courses.append(current_course)
                        current_course = {'overview': [], 'curriculum': []}
                    current_course['overview'].append(text)

                elif slide_type == "CURRICULUM":
                    current_course['curriculum'].append(text)

            if current_course['curriculum']:
                courses.append(current_course)

            print(f"  └─ 잠재 과정 수: {len(courses)}개")
            
            # ---------------------------------------------------------
            # LLM 변환 및 파일 저장
            # ---------------------------------------------------------
            for idx, course in enumerate(courses):
                full_overview = "\n\n".join(course['overview'])
                full_curriculum = "\n\n".join(course['curriculum'])
                
                md_content = generate_rag_markdown(file, idx+1, full_overview, full_curriculum)
                
                if md_content:
                    base_name = os.path.splitext(file)[0]
                    safe_name = re.sub(r'[^a-zA-Z0-9가-힣]', '_', base_name)
                    
                    md_filename = f"{safe_name}_Course_{idx+1}.md"
                    save_path = os.path.join(OUTPUT_DIR, md_filename)
                    
                    with open(save_path, 'w', encoding='utf-8') as f:
                        f.write(md_content)
                    
                    print(f"    ✅ Markdown 저장 완료: {md_filename}")
                else:
                    print(f"    🚫 [Drop] 과정 {idx+1}: 정보 부족")

        except Exception as e:
            print(f"  ❌ 파일 처리 중 에러 발생: {file} -> {e}")

    print(f"\n🎉 [커리큘럼] Markdown 변환 완료! '{OUTPUT_DIR}' 폴더를 확인하세요.")

if __name__ == "__main__":
    process_rag_data_final()